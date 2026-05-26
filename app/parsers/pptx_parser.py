from pptx import Presentation
from typing import List, Dict

class PPTXParser:
    def extract_text(self, file_path: str) -> List[Dict[str, str]]:
        """
        Extracts text from a PowerPoint presentation, preserving exact slide ordering.
        """
        extracted_slides = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                slide_content = []
                
                # Extract primary title
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    slide_content.append(f"Title: {slide.shapes.title.text.strip()}")
                
                # Extract all other text shapes
                for shape in slide.shapes:
                    if shape != slide.shapes.title and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_content.append(text)
                                
                # Crucially extract unwritten speaker context from notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"Notes: {notes}")
                
                extracted_slides.append({
                    "slide_number": slide_num + 1,
                    "content": "\n".join(slide_content)
                })
        except Exception as e:
            raise RuntimeError(f"Failed to deterministicly parse PPTX {file_path}: {e}")
            
        return extracted_slides
