from pptx import Presentation
import logging

logger = logging.getLogger("ElephantTank.Ingestion.PPTX")

class PPTXExtractor:
    @staticmethod
    def extract_text(file_path: str) -> str:
        """Extracts text from PPTX slides sequentially."""
        logger.info(f"Extracting PPTX text from: {file_path}")
        text_content = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {i+1} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text_content.append("\n".join(slide_text))
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Failed to parse PPTX {file_path}: {e}")
            raise ValueError(f"Could not parse PPTX: {e}")
